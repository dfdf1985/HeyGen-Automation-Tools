import os
import time
import subprocess
import requests
import re
import json
import traceback
from datetime import datetime, timedelta
from concurrent.futures import ThreadPoolExecutor
from dotenv import load_dotenv

# --- 第三方庫 ---
import qrcode
from googleapiclient.discovery import build  # 用於與 YouTube 官方 API 互動
from googleapiclient.errors import HttpError # 用於捕捉 YouTube API 的錯誤
import isodate  # 用於解析 YouTube 回傳的時間格式 (例如 PT5M30S)
from moviepy.editor import VideoFileClip, concatenate_videoclips # 用於影片剪輯與合併
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pdf2image import convert_from_path # 將 PDF 轉為圖片 (需安裝 poppler)

# ================= 1. 環境與參數設定 =================

# 載入 .env 檔案中的環境變數
load_dotenv()

# --- API Keys 檢查 ---
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
HEYGEN_API_KEY = os.getenv("HEYGEN_API_KEY")
YOUTUBE_API_KEY = os.getenv("YOUTUBE_API_KEY") 

# 確保所有必要的 API Key 都存在，否則拋出錯誤
if not GEMINI_API_KEY or not HEYGEN_API_KEY or not YOUTUBE_API_KEY:
    raise ValueError(" 錯誤：請確認 .env 檔案中包含 GEMINI, HEYGEN 和 YOUTUBE 的 API Key")

# --- 初始化 Clients ---
# Gemini Client (AI 文案生成)
client = genai.Client(api_key=GEMINI_API_KEY)

# YouTube Client (資料搜尋)
youtube_service = build('youtube', 'v3', developerKey=YOUTUBE_API_KEY)

# --- 用戶個人化設定 (可在此修改預設值) ---
USER_NAME = "j" 
CHANNEL_NAME = "科技全球焦點" 
TEMPLATE_PPPTX = "tech_template.pptx" # PPT 模板檔案名稱
OUTPUT_DIR = "outputs" # 輸出檔案的資料夾

# --- HeyGen API 參數 ---
API_HOST = "https://api.heygen.com"
GENERATE_URL_V2 = f"{API_HOST}/v2/video/generate" # 影片生成端點
UPLOAD_URL_V1 = "https://upload.heygen.com/v1/asset" # 資源上傳端點
VIDEO_STATUS_URL_V1 = f"{API_HOST}/v1/video_status.get" # 狀態查詢端點

# --- 數字人與聲音設定 (需替換為您自己的 HeyGen ID) ---
TALKING_PHOTO_ID = "8c6187262e744939bb335949024e3ec5" # 數字人像 ID
VOICE_ID_ZH = "4158cf2ef85d4ccc856aacb1c47dbb0c" # 中文語音 ID
VOICE_ID_EN = "cef3bc4e0a84424cafcde6f2cf466c97" # 英文語音 ID (備用)

# --- 語速設定 ---
# 1.0 = 原速, 1.1 = 稍快(推薦新聞感), 1.2 = 快, 0.9 = 慢
VOICE_SPEED = 1.1  

# --- 外部軟體路徑 ---
# LibreOffice 用於將 PPTX 轉為 PDF，請確保此路徑正確
WINDOWS_SOFFICE_PATH = r"C:\Program Files\LibreOffice\program\soffice.exe"

# ================= 2. 工具函數 =================

def safe_extract_json(text):
    """
    從 LLM (Gemini) 回傳的文字中提取 JSON 部分。
    處理包含 ```json ... ``` 標籤的情況。
    """
    if not text: return None
    match = re.search(r'```json\s*(.*?)\s*```', text, re.DOTALL)
    if match: return match.group(1)
    match = re.search(r'\[\s*\{.*\}\s*\]', text, re.DOTALL)
    if match: return match.group(0)
    return text.strip()

def clean_path_input(path_str):
    """
    清理使用者輸入的檔案路徑 (移除引號、& 符號等)。
    """
    if not path_str: return ""
    path_str = path_str.strip()
    if path_str.startswith("& "): path_str = path_str[2:].strip()
    path_str = path_str.strip('"').strip("'").strip('"').strip("'")
    return path_str

def generate_qr_code(url, output_path):
    """
    將網址轉換為 QR Code 圖片並存檔。
    """
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_L,
        box_size=10,
        border=2,
    )
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    img.save(output_path)
    return output_path

# --- [關鍵功能] 使用 YouTube Data API 搜尋影片 ---

def search_youtube_via_api(topic, limit=5):
    """
    使用官方 YouTube Data API 搜尋影片。
    流程：
    1. Search API: 找關鍵字相關的最新影片 ID
    2. Videos API: 根據 ID 查詢詳細資訊 (為了過濾 Shorts 和確認時間)
    """
    print(f" [API 搜尋] 正在搜尋：{topic} (啟用 Shorts 過濾)...")
    
    real_data = []
    
    try:
        # 步驟 1: 使用 Search.list 取得影片 ID
        # q=關鍵字, order='date' (最新的), type='video' (只找影片)
        search_response = youtube_service.search().list(
            q=topic,
            part='id',
            maxResults=limit * 3, # 多抓一點以便過濾掉 Shorts
            order='date', 
            type='video'
        ).execute()

        video_ids = [item['id']['videoId'] for item in search_response.get('items', [])]

        if not video_ids:
            print("    找不到任何相關影片 ID。")
            return []

        # 步驟 2: 使用 Videos.list 取得詳細資訊 (包含長度 contentDetails)
        videos_response = youtube_service.videos().list(
            id=','.join(video_ids),
            part='snippet,contentDetails'
        ).execute()

        count = 0
        for item in videos_response.get('items', []):
            if count >= limit: break

            # 解析影片長度 (格式如 PT5M30S)
            duration_iso = item['contentDetails']['duration']
            duration_obj = isodate.parse_duration(duration_iso)
            
            # 過濾 Shorts: 長度小於 60 秒視為 Shorts，跳過不處理
            if duration_obj.total_seconds() < 60:
                # print(f"    -> 跳過 Shorts: {item['snippet']['title']} ({duration_obj})")
                continue

            # 提取需要的資料
            title = item['snippet']['title']
            channel = item['snippet']['channelTitle']
            video_id = item['id']
            url = f"https://www.youtube.com/watch?v={video_id}"
            
            # 簡單清洗標題 (移除 HTML 實體符號，例如 &amp;)
            import html
            title = html.unescape(title)

            real_data.append({
                "title": title,
                "channel": channel,
                "url": url
            })
            
            print(f"    -> [V] 納入: {title[:15]}... ({str(duration_obj)})")
            count += 1
            
        return real_data

    except HttpError as e:
        print(f"    YouTube API Error: {e}")
        return []
    except Exception as e:
        print(f"    Search Error: {e}")
        return []

# ================= 3. 核心功能：內容生成 (API 搜尋 -> Gemini 改寫) =================

def fetch_content_and_make_pptx(topic, intro_script):
    """
    主邏輯：
    1. 呼叫 YouTube API 抓資料
    2. 呼叫 Gemini API 寫新聞稿
    3. 生成 PPTX 投影片 (含 QR Code)
    """
    now = datetime.now()
    current_date_str = now.strftime('%Y-%m-%d')
    
    # 步驟 1: 獲取真實資料
    real_videos = search_youtube_via_api(topic, limit=5)
    
    # 若 API 失敗或無資料，使用備用假資料 (避免程式崩潰)
    if not real_videos:
        print("    警告：搜尋不到任何影片，將使用備用資料。")
        real_videos = [{"title": "AI 趨勢分析", "channel": "科技頻道", "url": "https://www.youtube.com/"}]

    # 步驟 2: 將真實資料餵給 Gemini 進行改寫
    print(f" [1/5] 正在請 Gemini 撰寫播報稿...")
    
    videos_context = json.dumps(real_videos, ensure_ascii=False, indent=2)

    # Prompt 設計：要求 Gemini 扮演分析師，並輸出 JSON 格式
    prompt = f"""
    你是一位專業的 YouTube 趨勢分析師，今天是 {current_date_str}。
    我已經幫你搜尋好了以下 5 支與「{topic}」相關的真實 YouTube 影片資料：

    {videos_context}

    **任務**：
    請根據上述提供的【真實標題】與【頻道名稱】，為每一支影片撰寫介紹文案。
    **絕對不要修改網址 (url) 和 頻道名稱 (channel)**，直接沿用我給你的。

    **內容撰寫要求**：
    1. **highlight**: 請寫出一個 **「像新聞標題一樣簡潔」** 的標題，**嚴格限制在 20 個繁體中文字以內**。
    2. **speech**: 主播口播稿，格式：**「【頻道名稱】+【發布動作】+【核心內容】」**。
       - 範例：「【The Verge】發布最新評測，詳細比較了三款 AI 工具的優缺點。」

    **輸出格式 (JSON Array)**：
    [
        {{
            "channel": "原始資料的channel",
            "highlight": "簡潔標題",
            "speech": "口播稿",
            "url": "原始url" 
        }}
    ]
    """
    
    try:
        resp = client.models.generate_content(
            model="gemini-2.5-flash-lite", 
            contents=prompt,
            config=types.GenerateContentConfig(temperature=0.2)
        )
        
        full_text = ""
        if resp.candidates:
            for part in resp.candidates[0].content.parts:
                if hasattr(part, 'text') and part.text:
                    full_text += part.text
        
        json_str = safe_extract_json(full_text)
        final_data = json.loads(json_str)
        if not isinstance(final_data, list): final_data = [final_data]
        
    except Exception as e:
        print(f" Gemini 生成或解析錯誤: {e}")
        # Fallback: 如果 AI 失敗，使用原始資料簡單拼接
        final_data = []
        for v in real_videos:
            final_data.append({
                "channel": v['channel'],
                "highlight": v['title'][:20],
                "speech": f"【{v['channel']}】發布了關於 {v['title'][:10]} 的影片。",
                "url": v['url']
            })

    # --- 步驟 3: PPTX 生成 ---
    prs = Presentation(TEMPLATE_PPPTX) if os.path.exists(TEMPLATE_PPPTX) else Presentation()
    final_scripts = []
    all_collected_urls = [] 

    # 分頁邏輯：每頁放 3 則新聞
    ITEMS_PER_PAGE = 3 
    chunks = [final_data[i:i + ITEMS_PER_PAGE] for i in range(0, len(final_data), ITEMS_PER_PAGE)]

    for i, chunk in enumerate(chunks):
        # 選擇投影片版型 (如果有第二種版型則使用，否則用第一種)
        layout_idx = 1 if len(prs.slide_layouts) > 1 else 0
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx]) 
        
        if slide.shapes.title:
            slide.shapes.title.text = "熱門影片推薦"
        
        # 尋找文字框 (Placeholder)，若無則建立新的
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
        else:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
            tf = txBox.text_frame
        
        tf.clear()
        titles_on_page = [] 
        
        # 填入每一則新聞到 PPT
        for news_item in chunk:
            news_channel = news_item.get("channel", "YouTube")
            news_highlight = news_item.get("highlight", "")
            news_speech = news_item.get("speech", news_highlight) 
            news_url = news_item.get("url", "").strip()

            if not news_highlight: continue
            
            titles_on_page.append(news_speech)
            all_collected_urls.append(news_url)

            # [PPT 排版細節]
            p = tf.add_paragraph()
            p.space_after = Pt(20) 

            # 頻道名稱：金色、粗體、大字
            run_channel = p.add_run()
            run_channel.text = f"【{news_channel}】" 
            run_channel.font.size = Pt(24)
            run_channel.font.bold = True
            run_channel.font.color.rgb = RGBColor(255, 215, 0) # 金黃色

            # 新聞標題
            run_text = p.add_run()
            run_text.text = f" {news_highlight}"
            run_text.font.size = Pt(24)
            # run_text.font.color.rgb = RGBColor(255, 255, 255) # 若背景深色，可開啟此行

        # 組合口播稿
        if titles_on_page:
            titles_script = "。接著，".join(titles_on_page)
            if i == 0: 
                # 第一頁包含開場白
                final_str = f"{intro_script} 首先帶您關注：{titles_script}。"
            elif i == len(chunks) - 1: 
                # 最後一頁包含結尾語
                final_str = f"最後推薦：{titles_script}。以上是今天的 YouTube 趨勢整理。"
            else: 
                final_str = f"繼續看到：{titles_script}。"
            
            final_str = final_str.replace("。。", "。")
            final_scripts.append(final_str)
        else:
            final_scripts.append("以上是部分重點整理。")

    # --- QR Code 結尾頁面製作 ---
    if all_collected_urls:
        print("    >>> 正在製作 QR Code 結尾頁...")
        # 嘗試使用空白版型 (通常是 index 6)
        blank_layout_id = 6 if len(prs.slide_layouts) > 6 else 0
        last_slide = prs.slides.add_slide(prs.slide_layouts[blank_layout_id])
        
        # 加入標題
        title_box = last_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_tf = title_box.text_frame
        title_p = title_tf.add_paragraph()
        title_p.text = "影片來源列表 (掃描 QR Code 觀看)"
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 0, 0) 

        # 設定 QR Code 的起始位置與間距
        start_y = Inches(1.3)
        y_offset = Inches(1.15)
        qr_size = Inches(1.0)
        qr_x = Inches(0.8)
        text_x = Inches(2.0)
        
        temp_qr_files = [] 

        # 迴圈生成 QR Code 圖片並貼到 PPT 上 (最多顯示前 5 個)
        for i, url in enumerate(all_collected_urls[:5]):
            current_y = start_y + (i * y_offset)
            
            qr_filename = f"temp_qr_{i}.png"
            qr_path = os.path.join(OUTPUT_DIR, qr_filename)
            generate_qr_code(url, qr_path)
            temp_qr_files.append(qr_path)

            # 貼上圖片
            last_slide.shapes.add_picture(qr_path, qr_x, current_y, width=qr_size, height=qr_size)

            # 貼上文字連結
            url_box = last_slide.shapes.add_textbox(text_x, current_y + Inches(0.1), Inches(10.5), Inches(0.8))
            url_tf = url_box.text_frame
            url_tf.word_wrap = True
            
            url_p = url_tf.add_paragraph()
            display_text = url if len(url) < 80 else url[:77] + "..."
            
            run = url_p.add_run()
            run.text = f"[{i+1}] {display_text}"
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(5, 99, 193)
            run.font.underline = True
            try:
                run.hyperlink.address = url
            except ValueError: pass

        final_scripts.append("您可以掃描畫面上的 QR Code，或是點擊連結觀看完整影片。感謝您的收看，我們下次見！")

    # 儲存 PPTX
    pptx_path = os.path.join(OUTPUT_DIR, "final_gen.pptx")
    prs.save(pptx_path)

    # 清除暫存的 QR Code 圖片
    for qr_file in temp_qr_files:
        try: os.remove(qr_file)
        except: pass

    return pptx_path, final_scripts

# ================= 4. 圖片轉換與影片生成 =================

def convert_pptx_to_images(path):
    """
    將 PPTX 轉為圖片 (PNG) 序列，供 HeyGen 作為背景使用。
    流程：PPTX -> PDF (LibreOffice) -> PNGs (pdf2image)
    """
    print(" [2/5] PPT 轉圖片...")
    soffice = WINDOWS_SOFFICE_PATH if os.path.exists(WINDOWS_SOFFICE_PATH) else "soffice"
    
    # 呼叫 LibreOffice 轉檔指令
    try:
        subprocess.run([soffice, "--headless", "--convert-to", "pdf", path, "--outdir", OUTPUT_DIR], check=True)
    except FileNotFoundError:
        print(" 錯誤：找不到 LibreOffice，請確認路徑或安裝")
        raise

    # 將 PDF 轉為圖片
    pdf_path = os.path.join(OUTPUT_DIR, os.path.basename(path).replace(".pptx", ".pdf"))
    images = convert_from_path(pdf_path, dpi=200)
    paths = []
    for i, img in enumerate(images):
        p = os.path.join(OUTPUT_DIR, f"slide_{i+1}.png")
        img.save(p, "PNG")
        paths.append(p)
    return paths

def convert_custom_cover(file_path):
    """
    處理使用者上傳的自定義封面，轉為 PNG 格式。
    """
    ext = os.path.splitext(file_path)[1].lower()
    # 如果是 PPT 格式，先轉 PDF
    if ext in [".pptx", ".ppt"]:
        soffice = WINDOWS_SOFFICE_PATH if os.path.exists(WINDOWS_SOFFICE_PATH) else "soffice"
        subprocess.run([soffice, "--headless", "--convert-to", "pdf", file_path, "--outdir", OUTPUT_DIR], check=True)
        file_path = os.path.join(OUTPUT_DIR, os.path.basename(file_path).rsplit('.', 1)[0] + ".pdf")
        ext = ".pdf"
    # 如果是 PDF，取第一頁轉圖片
    if ext == ".pdf":
        images = convert_from_path(file_path, dpi=200, first_page=1, last_page=1)
        if images:
            save_path = os.path.join(OUTPUT_DIR, "custom_cover_final.png")
            images[0].save(save_path, "PNG")
            return save_path
    return file_path

# --- 上傳圖片到 HeyGen ---
def upload_to_heygen(file_path):
    """
    上傳單張圖片至 HeyGen，並取得 Asset ID。
    """
    headers = {"X-Api-Key": HEYGEN_API_KEY, "Content-Type": "image/png"}
    try:
        with open(file_path, "rb") as f: data = f.read()
        resp = requests.post(UPLOAD_URL_V1, headers=headers, data=data, params={"type": "image"})
        
        if not resp.ok:
            print(f"\n    [HeyGen 上傳失敗] 狀態碼: {resp.status_code}")
            print(f"    [錯誤回應]: {resp.text}")
            raise Exception(f"HeyGen 圖片上傳失敗: {resp.text}")
            
        return resp.json()["data"]["id"]
    except Exception as e:
        print(f"\n    [上傳異常]: {e}")
        raise

# --- 影片生成 ---
def create_full_video(image_paths, scripts):
    """
    發送請求給 HeyGen 生成影片。
    包含：背景圖、數字人 ID、講稿、語速。
    """
    print(f" [3/5] 生成 HeyGen 影片中...")
    # 使用多執行緒同時上傳多張投影片圖片，加速流程
    with ThreadPoolExecutor(max_workers=5) as executor:
        bg_ids = list(executor.map(upload_to_heygen, image_paths))
    
    scenes = []
    # 組合每一頁的場景 (Scene)
    for bg_id, script in zip(bg_ids, scripts):
        # 自動判斷語言 (如果有中文字就用中文語音，否則用英文)
        v_id = VOICE_ID_ZH if re.search(r"[\u4e00-\u9fff]", script) else VOICE_ID_EN
        scenes.append({
            "character": {
                "type": "talking_photo", 
                "talking_photo_id": TALKING_PHOTO_ID, 
                "scale": 0.25, # 數字人大小
                "offset": {"x": 0.4, "y": 0.4} # 數字人位置
            },
            "voice": {
                "type": "text", 
                "voice_id": v_id, 
                "input_text": script.replace('，', '， '), # 增加逗號停頓感
                "speed": VOICE_SPEED # 應用全域語速設定
            },
            "background": {
                "type": "image", 
                "image_asset_id": bg_id, 
                "fit": "contain"
            }
        })

    payload = {"video_inputs": scenes, "aspect_ratio": "16:9", "test": False, "caption": False}
    
    resp = requests.post(GENERATE_URL_V2, json=payload, headers={"X-Api-Key": HEYGEN_API_KEY})
    
    if not resp.ok:
        print(f"\n    [HeyGen 生成失敗] 狀態碼: {resp.status_code}")
        print(f"    [錯誤回應]: {resp.text}")
        raise Exception(f"HeyGen 影片生成請求被拒絕: {resp.text}")

    return resp.json()["data"]["video_id"]

def download_video(video_id, output_video_path):
    """
    輪詢 (Polling) HeyGen 狀態，直到影片渲染完成並下載。
    """
    print(" [4/5] 等待 HeyGen 渲染...")
    headers = {"X-Api-Key": HEYGEN_API_KEY}
    start_time = time.time()
    
    while True:
        try:
            r = requests.get(f"{VIDEO_STATUS_URL_V1}?video_id={video_id}", headers=headers).json()
            data = r.get("data", {})
            status = data.get("status")
        except: time.sleep(5); continue
        
        if status == "completed":
            print(f"\n    >>> HeyGen 渲染完成！下載中...")
            if data.get("video_url"):
                with open(output_video_path, "wb") as f: f.write(requests.get(data["video_url"]).content)
            break 
        elif status == "failed": raise Exception(f"渲染失敗: {data.get('error')}")
        
        # 顯示等待時間
        print(f"    ...已等待 {int(time.time()-start_time)} 秒 ({status})", end="\r")
        time.sleep(15) # 每 15 秒檢查一次

# ================= 5. 僅合併影片 (無字幕處理) =================

def merge_intro_and_news_video_only(intro_video_path, news_video_path, output_dir):
    """
    使用 MoviePy 將「自選片頭」與「HeyGen 生成的新聞影片」合併。
    """
    print(f" [5/5] 正在進行後製 (合併影片)...")
    merged_video_path = news_video_path

    try:
        clip_intro = VideoFileClip(intro_video_path)
        clip_news = VideoFileClip(news_video_path)
        print(f"    > 片頭長度: {clip_intro.duration:.2f} 秒")

        # 合併影片
        final_clip = concatenate_videoclips([clip_intro, clip_news], method="compose")
        merged_video_path = os.path.join(output_dir, "final_merged_output.mp4")
        
        # 輸出檔案
        final_clip.write_videofile(merged_video_path, codec="libx264", audio_codec="aac", verbose=False, logger=None)
        
        clip_intro.close()
        clip_news.close()
        
        print("    > 影片合併完成！")
        return merged_video_path

    except Exception as e:
        print(f"    >>> 合併錯誤: {e}")
        return news_video_path

# ================= 6. 主程式入口 =================

if __name__ == "__main__":
    try:
        os.makedirs(OUTPUT_DIR, exist_ok=True)
        print("\n=== AI YouTube 趨勢新聞生成器  ===\n")
        
        # 1. 獲取使用者輸入
        topic = input("1. YouTube 搜尋主題 (Enter 使用預設: 最新 AI 工具): ").strip() or "最新 AI 工具"
        default_intro = f"歡迎收看{CHANNEL_NAME}，我是{USER_NAME}。馬上帶您瀏覽 YouTube 最新熱點。"
        intro_script = input(f"2. 數字人開場白 (Enter 使用預設): ").strip() or default_intro
        
        raw_video_path = input(f"3. 片頭影片路徑 (請直接將檔案拖入視窗): ")
        custom_intro_video = clean_path_input(raw_video_path)
        
        if custom_intro_video and not os.path.exists(custom_intro_video):
            print(f"    警告: 找不到檔案，將忽略片頭。")
            custom_intro_video = ""

        custom_cover_img = "" 
        
        # 2. 爬取資料並製作 PPT
        pptx_path, scripts = fetch_content_and_make_pptx(topic, intro_script)
        
        # 3. PPT 轉圖片
        images = convert_pptx_to_images(pptx_path)

        # 處理自定義封面 (若有)
        if custom_cover_img and os.path.exists(custom_cover_img):
            final_cover_path = convert_custom_cover(custom_cover_img)
            if final_cover_path: images[0] = final_cover_path

        # 4. HeyGen 生成影片
        video_id = create_full_video(images, scripts)
        heygen_video_path = os.path.join(OUTPUT_DIR, "heygen_raw.mp4")
        download_video(video_id, heygen_video_path)
        
        # 5. 合併片頭 (若有)
        if custom_intro_video:
            final_vid = merge_intro_and_news_video_only(custom_intro_video, heygen_video_path, OUTPUT_DIR)
            print(f"\n    影片完成: {final_vid}")
            if os.name == 'nt': os.startfile(OUTPUT_DIR) # Windows 自動開啟資料夾
        else:
            print(f"\n！檔案位於: {heygen_video_path}")
            if os.name == 'nt': os.startfile(OUTPUT_DIR)
        
    except Exception as e:
        traceback.print_exc()
        input("錯誤，請按 Enter 離開...")