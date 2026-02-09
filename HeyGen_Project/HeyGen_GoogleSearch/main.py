import os
import time
import subprocess
import requests
import re
import json
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor # 用於多執行緒加速上傳
from dotenv import load_dotenv

# --- SDK 與 工具 導入 ---
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Pt, Inches
from pdf2image import convert_from_path

# ================= 1. 環境與參數設定 =================

# 載入 .env 檔案中的環境變數 (保護 API Key 不外洩)
load_dotenv()

# --- API Keys 檢查 (防呆機制) ---
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
HEYGEN_API_KEY = os.getenv("HEYGEN_API_KEY")

if not GEMINI_API_KEY or not HEYGEN_API_KEY:
    raise ValueError(" 錯誤：請確認 .env 檔案中包含有效的 API Key")

# 初始化 Google Gemini Client
client = genai.Client(api_key=GEMINI_API_KEY)

# --- 用戶個人化設定 ---
USER_NAME = "j" 
CHANNEL_NAME = "科技全球焦點" 
TEMPLATE_PPPTX = "tech_template.pptx" # 簡報模板檔案名稱
OUTPUT_DIR = "outputs"                # 輸出檔案存放目錄
FINAL_VIDEO_NAME = "final_news_video.mp4"

# --- HeyGen API 參數 ---
API_HOST = "https://api.heygen.com"
GENERATE_URL_V2 = f"{API_HOST}/v2/video/generate" # v2 影片生成接口
UPLOAD_URL_V1 = "https://upload.heygen.com/v1/asset" # v1 資產上傳接口
VIDEO_STATUS_URL_V1 = f"{API_HOST}/v1/video_status.get" # 查詢生成狀態

# --- 數字人 (Avatar) 與 聲音 ID 設定 ---
# 建議：這些 ID 可以從 HeyGen 網頁版 URL 或 API 列表獲取
TALKING_PHOTO_ID = "8c6187262e744939bb335949024e3ec5"
VOICE_ID_ZH = "4158cf2ef85d4ccc856aacb1c47dbb0c" # 中文聲音
VOICE_ID_EN = "cef3bc4e0a84424cafcde6f2cf466c97" # 英文聲音 (備用)

# ⚠️ LibreOffice 路徑設定
# 這是將 PPT 轉為 PDF 的關鍵工具，請確保路徑與您電腦安裝位置一致
WINDOWS_SOFFICE_PATH = r"C:\Program Files\LibreOffice\program\soffice.exe"

# ================= 2. 工具函數 =================

def safe_extract_json(text):
    """
    從 AI 回傳的文字中安全提取 JSON 字串。
    AI 有時會包裹 Markdown 標籤 (```json ... ```)，此函數用於去除這些雜訊。
    """
    if not text: return None
    # 嘗試抓取 ```json 包裹的內容
    match = re.search(r'```json\s*(.*?)\s*```', text, re.DOTALL)
    if match: return match.group(1)
    # 若無 Markdown，嘗試抓取最外層的陣列 []
    match = re.search(r'\[\s*\{.*\}\s*\]', text, re.DOTALL)
    if match: return match.group(0)
    return text.strip()

# ================= 3. 核心功能：內容生成 =================

def fetch_content_and_make_pptx(topic, intro_script):
    """
    流程 A: 使用 Gemini 聯網搜尋新聞 -> 整理成 JSON -> 製作 PPT -> 生成口播稿
    """
    now = datetime.now()
    current_date_str = now.strftime('%Y-%m-%d')
    
    print(f" [1/5] 正在搜尋「{topic}」")
    
    # --- A. Gemini 搜尋與內容生成 ---
    # Prompt 設計重點：指定角色、時間、強制 JSON 格式、限制字數
    prompt = f"""
    你是一位資深新聞編輯，今天是 {current_date_str}。
    
    任務目標：使用 Google Search 工具，針對主題「{topic}」搜尋最近 24 小時至本週內的重大事件。
    
    請嚴格遵守以下規則：
    1. **強制搜尋與去重**：務必執行搜尋。若有多家媒體報導同一事件，請合併為單一條目。
    2. **數量限制**：請精選出 **5 則** 最具影響力的新聞。
    3. **內容撰寫**：
       - 請將每則新聞濃縮為一句 **「約 20 個中文字」** 的重點解說。
       - 包含「主詞」+「發生什麼事」+「結果」。
       - 範例：「輝達財報優於預期，AI晶片需求強勁，盤後股價大漲10%。」
    4. **輸出格式**：
       嚴格以 JSON Array 格式輸出，key 使用 "summary"：
       [
         {{ "summary": "新聞1的20字解說..." }}
       ]
    """
    
    # 設定工具：啟用 Google Search
    search_tool = types.Tool(google_search=types.GoogleSearch())

    try:
        # 呼叫 Gemini API
        resp = client.models.generate_content(
            model="gemini-2.5-flash-lite", # 使用輕量快速的模型
            contents=prompt,
            config=types.GenerateContentConfig(tools=[search_tool], temperature=0.2) # 低溫創造性，求精準
        )
        
        full_text = ""
        if resp.candidates:
            for part in resp.candidates[0].content.parts:
                if hasattr(part, 'text') and part.text:
                    full_text += part.text
        
        # 解析 JSON 資料
        raw_data = json.loads(safe_extract_json(full_text))
        if not isinstance(raw_data, list): raw_data = [raw_data] # 確保格式是 List
        
        print(f"   >>> 搜尋完成，第一筆資料範例: {raw_data[0] if raw_data else '無資料'}")
        print(f"   >>> 共生成 {len(raw_data)} 則新聞。")
            
    except Exception as e:
        print(f" Gemini 錯誤: {e}")
        # 錯誤處理：若 API 失敗，生成一條假資料讓程式能繼續跑，方便 Debug
        raw_data = [{"summary": f"今日{topic}相關新聞整理 (擷取失敗，請檢查 API 或網絡)"}]

    # --- B. PPTX 生成 ---
    
    # 載入模板或建立新簡報
    prs = Presentation(TEMPLATE_PPPTX) if os.path.exists(TEMPLATE_PPPTX) else Presentation()
    final_scripts = []

    # 🟢 Slide 1: 封面製作
    cover_idx = 0 if len(prs.slide_layouts) > 0 else 0
    slide_cover = prs.slides.add_slide(prs.slide_layouts[cover_idx]) 
    
    # 設定標題與副標題
    if slide_cover.shapes.title:
        slide_cover.shapes.title.text = f"{topic} 快報"
    if len(slide_cover.placeholders) > 1:
        slide_cover.placeholders[1].text = f"{CHANNEL_NAME} | {current_date_str}"
    
    # 加入開場白 script
    final_scripts.append(intro_script)

    # 🟢 Slide 2~N: 分頁內容 (分頁邏輯)
    ITEMS_PER_PAGE = 3 # 每頁最多顯示 3 則新聞
    # 將新聞列表切塊 (Chunking)，例如 5 則新聞會變成 [3則, 2則] 兩頁
    chunks = [raw_data[i:i + ITEMS_PER_PAGE] for i in range(0, len(raw_data), ITEMS_PER_PAGE)]

    for i, chunk in enumerate(chunks):
        layout_idx = 1 if len(prs.slide_layouts) > 1 else 0
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx]) 
        
        # 設定分頁標題 (第一頁顯示「摘要」，後續顯示「摘要(續)」)
        if slide.shapes.title:
            slide.shapes.title.text = "今日重點摘要" if i == 0 else "今日重點摘要 (續)"
        
        # 取得文字框 (優先使用佔位符，若無則建立新文字框)
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
        else:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
            tf = txBox.text_frame
        
        tf.clear() # 清空預設文字
        
        # --- 構建該頁面的內容與口播稿 ---
        titles_on_page = []
        
        for news_item in chunk:
            # 優先讀取 "summary"，若無則讀取 "title"
            news_content = news_item.get("summary", news_item.get("title", ""))
            
            if not news_content: continue

            titles_on_page.append(news_content)

            # 將新聞寫入 PPT
            p = tf.add_paragraph()
            p.text = f"{news_content}" 
            p.font.size = Pt(24) 
            p.font.bold = True
            p.space_after = Pt(24) # 段落間距

        # --- 自動生成流暢的過場口播 ---
        if titles_on_page:
            titles_script = "。接著是，".join(titles_on_page)
            
            # 根據是第幾頁，調整連接詞 (首先、繼續關注、最後)
            if len(chunks) == 1:
                final_str = f"今天的重點包括：{titles_script}。以上是今天的快報，感謝收看。"
            else:
                if i == 0:
                    final_str = f"首先帶您關注：{titles_script}。"
                elif i == len(chunks) - 1:
                    final_str = f"最後看到：{titles_script}。以上是今天的快報，感謝收看。"
                else:
                    final_str = f"繼續關注：{titles_script}。"
            
            # 修正標點符號，避免多個句號
            final_str = final_str.replace("。。", "。")
            final_scripts.append(final_str)
        else:
            final_scripts.append("以上是部分重點整理。")

    pptx_path = os.path.join(OUTPUT_DIR, "final_gen.pptx")
    prs.save(pptx_path)
    return pptx_path, final_scripts

# ================= 4. 圖片轉換與影片生成 =================

def convert_pptx_to_images(path):
    """
    使用 LibreOffice 將 PPTX -> PDF -> PNG 圖片
    """
    print(" [2/5] PPT 轉圖片...")
    soffice = WINDOWS_SOFFICE_PATH if os.path.exists(WINDOWS_SOFFICE_PATH) else "soffice"
    
    # 呼叫系統指令執行轉檔 (--headless 代表不開啟圖形介面，背景執行)
    subprocess.run([soffice, "--headless", "--convert-to", "pdf", path, "--outdir", OUTPUT_DIR], check=True)
    
    pdf_path = os.path.join(OUTPUT_DIR, os.path.basename(path).replace(".pptx", ".pdf"))
    
    # 將 PDF 每一頁轉為圖片
    images = convert_from_path(pdf_path, dpi=200)
    paths = []
    for i, img in enumerate(images):
        p = os.path.join(OUTPUT_DIR, f"slide_{i+1}.png")
        img.save(p, "PNG")
        paths.append(p)
    return paths

def convert_custom_cover(file_path):
    """
    處理使用者上傳的自訂封面 (支援 PPTX 或 PDF 轉為 PNG)
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    # 若是 PPT 格式，先轉 PDF
    if ext in [".pptx", ".ppt"]:
        print(f"處理自訂 PPT 封面...")
        soffice = WINDOWS_SOFFICE_PATH if os.path.exists(WINDOWS_SOFFICE_PATH) else "soffice"
        subprocess.run([soffice, "--headless", "--convert-to", "pdf", file_path, "--outdir", OUTPUT_DIR], check=True)
        file_path = os.path.join(OUTPUT_DIR, os.path.basename(file_path).rsplit('.', 1)[0] + ".pdf")
        ext = ".pdf"
        
    # 若是 PDF 格式，取第一頁轉圖片
    if ext == ".pdf":
        print(f"處理 PDF 封面...")
        images = convert_from_path(file_path, dpi=200, first_page=1, last_page=1)
        if images:
            save_path = os.path.join(OUTPUT_DIR, "custom_cover_final.png")
            images[0].save(save_path, "PNG")
            return save_path
            
    return file_path # 如果原本就是圖片，直接回傳路徑

def upload_to_heygen(file_path):
    """
    將單張圖片上傳至 HeyGen 資產庫，取得 asset_id
    """
    headers = {"X-Api-Key": HEYGEN_API_KEY, "Content-Type": "image/png"}
    with open(file_path, "rb") as f: data = f.read()
    # 參數 type=image 告訴 HeyGen 這是圖片
    resp = requests.post(UPLOAD_URL_V1, headers=headers, data=data, params={"type": "image"})
    return resp.json()["data"]["id"]

def create_full_video(image_paths, scripts):
    """
    組合 HeyGen 影片場景
    """
    print(f" [3/5] 生成影片中 (素材上傳與合成)...")
    
    # 使用 ThreadPool 平行上傳圖片，加快速度 (一次上傳 5 張)
    with ThreadPoolExecutor(max_workers=5) as executor:
        bg_ids = list(executor.map(upload_to_heygen, image_paths))
    
    scenes = []
    # 將每一張圖片 (bg_id) 與對應的口播稿 (script) 配對
    for bg_id, script in zip(bg_ids, scripts):
        # 簡單的語言判斷：如果有中文字就用中文語音，否則用英文
        v_id = VOICE_ID_ZH if re.search(r"[\u4e00-\u9fff]", script) else VOICE_ID_EN
        
        scenes.append({
            "character": {
                "type": "talking_photo", 
                "talking_photo_id": TALKING_PHOTO_ID, 
                "scale": 0.25, # 頭像大小
                "offset": {"x": 0.4, "y": 0.4} # 頭像位置
            },
            "voice": {
                "type": "text", 
                "voice_id": v_id, 
                "input_text": script.replace('，', '， ') # 增加逗號後的停頓感
            },
            "background": {
                "type": "image", 
                "image_asset_id": bg_id, 
                "fit": "contain" # 背景適應方式
            }
        })

    # 發送生成請求
    payload = {"video_inputs": scenes, "aspect_ratio": "16:9", "test": False, "caption": True}
    resp = requests.post(GENERATE_URL_V2, json=payload, headers={"X-Api-Key": HEYGEN_API_KEY})
    return resp.json()["data"]["video_id"]

def download_video(video_id, output_video_path):
    """
    輪詢 (Polling) 檢查影片生成狀態，完成後下載
    """
    print(" [4/5] 等待 HeyGen 渲染...")
    headers = {"X-Api-Key": HEYGEN_API_KEY}
    start_time = time.time()
    
    while True:
        try:
            r = requests.get(f"{VIDEO_STATUS_URL_V1}?video_id={video_id}", headers=headers).json()
            data = r.get("data", {})
            status = data.get("status")
        except: 
            time.sleep(5); continue # 若網路請求失敗，稍微等待重試
        
        if status == "completed":
            print(f"\n   >>> 渲染完成！下載中...")
            # 下載影片
            if data.get("video_url"):
                with open(output_video_path, "wb") as f: f.write(requests.get(data["video_url"]).content)
            # 下載字幕 (如果有)
            if data.get("caption_url"):
                with open(output_video_path.replace(".mp4", ".srt"), "wb") as f: f.write(requests.get(data["caption_url"]).content)
            break
        elif status == "failed": 
            raise Exception(f"渲染失敗: {data.get('error')}")
        
        # 顯示等待秒數
        print(f"   ...已等待 {int(time.time()-start_time)} 秒 ({status})", end="\r")
        time.sleep(15) # 每 15 秒檢查一次，避免過於頻繁

# ================= 5. 主程式入口 =================

if __name__ == "__main__":
    try:
        # 確保輸出目錄存在
        os.makedirs(OUTPUT_DIR, exist_ok=True)
        
        # 使用者輸入互動
        topic = input("1. 搜尋主題: (例如: AI最新趨勢):").strip() or "AI最新趨勢"
        default_intro = f"歡迎收看{CHANNEL_NAME}，我是{USER_NAME}。馬上帶您瀏覽重點標題。"
        intro_script = input(f"2. 開場白 (Enter 使用預設): ").strip() or default_intro
        
        # 處理路徑字串中可能包含的引號
        custom_cover = input(f"3. 封面路徑 (選填): ").strip().replace('"', '').replace("'", "")

        # 步驟 1: 獲取內容並製作 PPT
        pptx_path, scripts = fetch_content_and_make_pptx(topic, intro_script)
        
        # 步驟 2: 轉換圖片
        images = convert_pptx_to_images(pptx_path)

        # 步驟 3: 若有自訂封面，替換第一張圖片
        if custom_cover and os.path.exists(custom_cover):
            final_cover_path = convert_custom_cover(custom_cover)
            if final_cover_path: images[0] = final_cover_path

        # 步驟 4: 上傳並生成影片
        video_id = create_full_video(images, scripts)
        
        # 步驟 5: 下載成品
        download_video(video_id, os.path.join(OUTPUT_DIR, FINAL_VIDEO_NAME))
        
        print(f"\n  完成！影片位置: {os.path.join(OUTPUT_DIR, FINAL_VIDEO_NAME)}")
        
    except Exception as e:
        import traceback
        traceback.print_exc() # 印出詳細錯誤訊息以便除錯
        input("錯誤，請按 Enter 離開...")